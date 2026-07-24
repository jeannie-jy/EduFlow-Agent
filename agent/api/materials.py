"""素材 API 路由。

POST   /api/materials/upload                      上传课件文件
POST   /api/materials/{id}/parse                  解析为结构化内容
GET    /api/materials/{id}/preview                预览解析结果
"""

from __future__ import annotations

import logging
import shutil
import uuid
from pathlib import Path

from fastapi import APIRouter, Depends, HTTPException, UploadFile, File
from sqlalchemy import select
from sqlalchemy.ext.asyncio import AsyncSession

from config import get_settings
from db.database import get_session
from db.models import SourceMaterial
from .deps import CurrentUser

logger = logging.getLogger(__name__)

router = APIRouter(prefix="/materials", tags=["materials"])

# 文件类型白名单（扩展名 → MIME）
ALLOWED_EXTENSIONS = {
    ".pdf": "application/pdf",
    ".txt": "text/plain",
    ".md": "text/markdown",
    ".py": "text/x-python",
    ".c": "text/x-csrc",
    ".java": "text/x-java-source",
    ".cpp": "text/x-c++src",
    ".pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
}


def _normalize_original_filename(filename: str | None) -> str:
    """Return a database-safe display name without changing the stored path."""
    basename = (filename or "unknown").replace("\\", "/").rsplit("/", 1)[-1]
    cleaned = "".join(
        character for character in basename if ord(character) >= 32 and ord(character) != 127
    ).strip()
    if not cleaned:
        return "unknown"

    suffix = Path(cleaned).suffix
    if len(cleaned) > 500:
        if suffix and len(suffix) < 500:
            return f"{cleaned[:500 - len(suffix)]}{suffix}"
        return cleaned[:500]
    return cleaned


def resolve_material_storage_path(storage_path: str | None) -> Path:
    """Resolve a persisted path only when it remains inside the upload root."""
    if storage_path is None:
        raise HTTPException(status_code=404, detail="Material not found")

    upload_root = get_settings().upload_dir.resolve()
    try:
        file_path = Path(storage_path).resolve()
        file_path.relative_to(upload_root)
    except (OSError, ValueError):
        raise HTTPException(status_code=404, detail="Material not found")

    if not file_path.is_file():
        raise HTTPException(status_code=404, detail="Material not found")
    return file_path


@router.post("/upload", status_code=201)
async def upload_material(
    current_user: CurrentUser,
    file: UploadFile = File(...),
    session: AsyncSession = Depends(get_session),
) -> dict:
    """上传课件文件。

    支持: PDF, PPTX, Markdown, TXT, Python/C/Java/C++ 代码。
    最大: 50 MB。
    """
    settings = get_settings()

    # 文件名安全检查
    original_name = _normalize_original_filename(file.filename)
    suffix = Path(original_name).suffix.lower()

    if suffix not in ALLOWED_EXTENSIONS:
        raise HTTPException(
            status_code=400,
            detail=f"不支持的文件类型: {suffix}。支持: {', '.join(ALLOWED_EXTENSIONS.keys())}",
        )

    expected_content_type = ALLOWED_EXTENSIONS[suffix]
    if (
        expected_content_type not in settings.allowed_upload_types
        or file.content_type != expected_content_type
    ):
        raise HTTPException(status_code=400, detail="File content type does not match its extension")

    # 大小检查
    contents = await file.read()
    if len(contents) > settings.upload_max_size_bytes:
        raise HTTPException(
            status_code=413,
            detail=f"文件过大（最大 {settings.upload_max_size_bytes // 1048576} MB）",
        )

    material_id = uuid.uuid4()
    upload_dir = settings.upload_dir / str(material_id)
    upload_dir.mkdir(parents=True, exist_ok=True)

    # 保存文件
    safe_filename = f"uploaded_{material_id.hex[:8]}{suffix}"
    file_path = upload_dir / safe_filename
    file_path.write_bytes(contents)

    material = SourceMaterial(
        id=material_id,
        owner_id=current_user.id,
        project_id=None,
        type=suffix.lstrip("."),
        filename=original_name,
        size_bytes=len(contents),
        storage_path=str(file_path),
    )
    try:
        session.add(material)
        await session.flush()
        await session.commit()
    except BaseException:
        # Cleanup is best-effort so it never masks the original database error
        # or cancellation signal.
        try:
            await session.rollback()
        except BaseException:
            logger.exception("素材上传失败后的数据库回滚失败: id=%s", material_id)
        try:
            shutil.rmtree(upload_dir)
        except OSError:
            logger.exception("素材上传失败后的文件清理失败: id=%s", material_id)
        raise

    logger.info("文件上传: id=%s | name=%s | size=%d | type=%s",
                material_id, original_name, len(contents), suffix)

    return {
        "id": str(material_id),
        "filename": original_name,
        "type": suffix.lstrip("."),
        "size_bytes": len(contents),
    }


def parse_material_file(storage_path: str | Path) -> dict | None:
    """Parse a stored material path into ``{topics, raw_text}``.

    Callers must resolve the path from an authorized database record first.
    """
    file_path = Path(storage_path)
    if not file_path.is_file():
        return None

    suffix = file_path.suffix.lower()

    raw_text = ""
    topics: list[str] = []
    if suffix == ".pdf":
        raw_text, topics = _parse_pdf(file_path)
    elif suffix == ".pptx":
        raw_text, topics = _parse_pptx(file_path)
    elif suffix in (".txt", ".md", ".py", ".c", ".java", ".cpp"):
        raw_text = file_path.read_text(encoding="utf-8")
        topics = _extract_topics_from_text(raw_text)
    else:
        raw_text = f"Unsupported format: {suffix}"

    return {"topics": topics, "raw_text": raw_text[:100000]}  # 截断到 100K 字符


async def get_owned_material(
    session: AsyncSession,
    material_id: str,
    user_id: uuid.UUID,
) -> SourceMaterial:
    """Return a material only when it belongs to the authenticated user."""
    try:
        material_uuid = uuid.UUID(material_id)
    except (ValueError, AttributeError):
        raise HTTPException(status_code=404, detail="Material not found")

    material = await session.scalar(
        select(SourceMaterial).where(
            SourceMaterial.id == material_uuid,
            SourceMaterial.owner_id == user_id,
        )
    )
    if material is None:
        raise HTTPException(status_code=404, detail="Material not found")
    return material


async def get_owned_project_material(
    session: AsyncSession,
    material_id: str,
    user_id: uuid.UUID,
    project_id: uuid.UUID,
) -> SourceMaterial:
    """Return a material only when it belongs to both user and project."""
    try:
        material_uuid = uuid.UUID(material_id)
    except (ValueError, AttributeError):
        raise HTTPException(status_code=404, detail="Material not found")

    material = await session.scalar(
        select(SourceMaterial).where(
            SourceMaterial.id == material_uuid,
            SourceMaterial.owner_id == user_id,
            SourceMaterial.project_id == project_id,
        )
    )
    if material is None:
        raise HTTPException(status_code=404, detail="Material not found")
    return material


@router.post("/{material_id}/parse")
async def parse_material(
    material_id: str,
    current_user: CurrentUser,
    session: AsyncSession = Depends(get_session),
) -> dict:
    """解析上传的课件文件，提取结构化内容。"""
    material = await get_owned_material(session, material_id, current_user.id)
    file_path = resolve_material_storage_path(material.storage_path)

    try:
        parsed = parse_material_file(file_path)
    except Exception as exc:
        logger.exception("文件解析失败: material=%s", material_id)
        raise HTTPException(status_code=500, detail=f"解析失败: {exc}")

    if parsed is None:
        raise HTTPException(status_code=404, detail="Material not found")

    material.content_text = parsed["raw_text"]
    material.parsed_result = parsed
    await session.flush()

    logger.info("文件解析完成: id=%s | topics=%d | text_len=%d",
                material_id, len(parsed["topics"]), len(parsed["raw_text"]))

    return {
        "id": str(material.id),
        "status": "done",
        "parsed_result": parsed,
    }


@router.get("/{material_id}/preview")
async def preview_material(
    material_id: str,
    current_user: CurrentUser,
    session: AsyncSession = Depends(get_session),
) -> dict:
    """预览文件解析结果（只读，不触发重新解析）。"""
    material = await get_owned_material(session, material_id, current_user.id)
    file_path = resolve_material_storage_path(material.storage_path)

    suffix = file_path.suffix.lower()

    preview_text = ""
    if suffix in (".txt", ".md", ".py", ".c", ".java", ".cpp"):
        preview_text = file_path.read_text(encoding="utf-8")[:5000]
    else:
        preview_text = f"Binary file: {file_path.name} ({suffix})"

    return {
        "id": str(material.id),
        "filename": material.filename,
        "type": material.type,
        "size_bytes": material.size_bytes,
        "preview": preview_text,
    }


# ============================================================================
# 解析器实现
# ============================================================================


def _parse_pdf(file_path: Path) -> tuple[str, list[str]]:
    """解析 PDF 文件，提取文本和候选主题。"""
    try:
        import fitz  # PyMuPDF
    except ImportError:
        return ("PDF 解析器未安装 (pip install PyMuPDF)", [])

    doc = fitz.open(str(file_path))
    pages_text = []
    for page in doc:
        pages_text.append(page.get_text())

    full_text = "\n".join(pages_text)
    topics = _extract_topics_from_text(full_text)
    doc.close()
    return full_text, topics


def _parse_pptx(file_path: Path) -> tuple[str, list[str]]:
    """解析 PPTX 文件。"""
    try:
        from pptx import Presentation
    except ImportError:
        return ("PPTX 解析器未安装 (pip install python-pptx)", [])

    prs = Presentation(str(file_path))
    slides_text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                slides_text.append(shape.text_frame.text)

    full_text = "\n---\n".join(slides_text)
    topics = _extract_topics_from_text(full_text)
    return full_text, topics


def _extract_topics_from_text(text: str) -> list[str]:
    """从文本中提取候选知识主题（基于 CS 关键词匹配）。

    MVP: 简单的关键词密度匹配。
    完整版: 使用 LLM 的 extract_concepts Tool。
    """
    if not text.strip():
        return []

    # CS 关键词库
    cs_keywords = [
        # 数据结构
        "数组", "链表", "栈", "队列", "哈希表", "堆", "二叉树", "AVL",
        "红黑树", "B树", "B+树", "图", "邻接表", "邻接矩阵",
        # 算法
        "排序", "冒泡", "快速排序", "归并排序", "二分", "递归", "动态规划",
        "贪心", "BFS", "DFS", "Dijkstra", "最小生成树", "拓扑排序",
        "最短路径", "松弛操作",
        # 操作系统
        "进程", "线程", "调度", "同步", "互斥", "死锁", "分页", "缓存",
        "虚拟内存", "TLB", "缺页", "信号量", "管程",
        # 网络
        "TCP", "UDP", "IP", "HTTP", "DNS", "路由", "拥塞控制", "三次握手",
        "四次挥手", "OSI", "子网",
        # 数据库
        "索引", "事务", "锁", "隔离级别", "B+树", "查询优化", "ACID",
        "连接", "SQL",
        # 软件工程
        "设计模式", "架构", "微服务", "CI/CD", "测试", "敏捷",
    ]

    # 计算每个关键词在文本中出现的次数
    scored = []
    for kw in cs_keywords:
        count = text.count(kw)
        if count > 0:
            scored.append((kw, count))

    # 按出现频率排序，取前 10
    scored.sort(key=lambda x: x[1], reverse=True)
    return [s[0] for s in scored[:10]]
